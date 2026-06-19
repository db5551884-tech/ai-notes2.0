import os
import io

from flask import Flask, render_template, request, send_file
from dotenv import load_dotenv

import google.generativeai as genai

import pdfplumber
from docx import Document
from pptx import Presentation
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet

# ========================
# INIT
# ========================

load_dotenv()

genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
model = genai.GenerativeModel("gemini-1.5-flash")

app = Flask(__name__)

lecture_text = ""


# ========================
# AI FUNCTIONS
# ========================

def summarize(text):
    if not text:
        return "Немає тексту для аналізу"

    try:
        prompt = f"""
Зроби:

1. Короткий конспект у вигляді пунктів
2. 5 тестових питань по тексту

Текст:
{text[:3000]}
"""
        response = model.generate_content(prompt)
        return response.text.replace("\n", "<br>")

    except Exception as e:
        return f"Помилка AI: {str(e)}"


def ask_ai(question, context):
    try:
        prompt = f"""
Відповідай на питання по тексту:

{context[:2000]}

Питання: {question}
"""
        response = model.generate_content(prompt)
        return response.text

    except Exception as e:
        return f"Помилка AI: {str(e)}"


# ========================
# FILE PROCESSING
# ========================

def extract_text(file):
    filename = file.filename.lower()

    if filename.endswith(".pdf"):
        text = ""
        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
        return text

    elif filename.endswith(".docx"):
        doc = Document(file)
        return "\n".join([p.text for p in doc.paragraphs])

    elif filename.endswith(".pptx"):
        ppt = Presentation(file)
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    return ""


# ========================
# ROUTES
# ========================

@app.route("/", methods=["GET", "POST"])
def index():
    global lecture_text

    summary = ""
    answer = ""

    if request.method == "POST":
        text = request.form.get("text")
        file = request.files.get("file")

        if file and file.filename != "":
            lecture_text = extract_text(file)
            summary = summarize(lecture_text)

        elif text:
            lecture_text = text
            summary = summarize(text)

    return render_template("index.html", summary=summary, answer=answer)


@app.route("/ask", methods=["POST"])
def ask():
    global lecture_text

    question = request.form.get("question")

    if lecture_text:
        answer = ask_ai(question, lecture_text[:2000])
    else:
        answer = "Спочатку введи текст або завантаж файл"

    return render_template("index.html", summary="", answer=answer)


# ========================
# PDF DOWNLOAD
# ========================

def create_pdf(text):
    buffer = io.BytesIO()

    doc = SimpleDocTemplate(buffer)
    styles = getSampleStyleSheet()

    content = []

    for line in text.split("<br>"):
        content.append(Paragraph(line, styles["Normal"]))

    doc.build(content)

    buffer.seek(0)
    return buffer


@app.route("/download-pdf")
def download_pdf():
    summary = request.args.get("text", "")

    if not summary:
        return "Немає тексту для PDF"

    pdf = create_pdf(summary)

    return send_file(
        pdf,
        as_attachment=True,
        download_name="summary.pdf",
        mimetype="application/pdf"
    )


# ========================
# RUN (Render friendly)
# ========================

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 10000))
    app.run(host="0.0.0.0", port=port)